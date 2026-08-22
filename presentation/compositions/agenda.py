"""
presentation/compositions/agenda.py
────────────────────────────────────────────────────────────────────────────
AGENDA composition — numbered table-of-contents slide.

Visual structure (inspired by Presenton's content_index_slide):

  ┌──────────────────────────────────────────────────────────┐
  │  SLIDE TITLE                                 [accent bar] │
  │──────────────────────────────────────────────────────────│
  │                                                           │
  │  01  ──────────────  02  ──────────────                  │
  │      Item title          Item title                      │
  │      Item subtitle       Item subtitle                   │
  │                                                           │
  │  03  ──────────────  04  ──────────────                  │
  │      Item title          Item title                      │
  │      Item subtitle       Item subtitle                   │
  │                                                           │
  │  05  ──────────────                                      │
  │      Item title                                          │
  │      Item subtitle                                       │
  │                                                           │
  └──────────────────────────────────────────────────────────┘

Each item:
  - Large number (accent color, bold, ~52pt)
  - Thin accent separator line beneath number
  - Title  (primary color, bold, ~15pt)
  - Subtitle (text_dark, regular, ~12pt)

Layout adapts to item count:
  1–2   → single row, items centered
  3–4   → 2 columns × up to 2 rows
  5–6   → 2 columns × up to 3 rows
  7–9   → 3 columns × up to 3 rows

Public API:
    render_agenda(pptx_slide, slide: SlideData, theme: Theme) -> None
"""

from __future__ import annotations

import logging
import math
from typing import Any

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ai.schemas import SlideData
from presentation.styles import Theme, DesignTokens

logger = logging.getLogger(__name__)

TK = DesignTokens()

# ── slide dimensions ──────────────────────────────────────────────────────
SLIDE_W = 13.33   # inches
SLIDE_H = 7.5

# ── agenda-specific tokens ────────────────────────────────────────────────
NUMBER_SIZE       = 52    # pt  — large bold ordinal
TITLE_SIZE        = 15    # pt  — item heading
SUBTITLE_SIZE     = 11    # pt  — item description
NUMBER_H          = 0.75  # in  — height of number textbox
SEP_H             = 0.04  # in  — thin separator line under number
GAP_SEP_TITLE     = 0.10  # in  — gap between separator and title
TITLE_H           = 0.32  # in  — title textbox height
SUBTITLE_H        = 0.30  # in  — subtitle textbox height
ITEM_BOTTOM_PAD   = 0.20  # in  — extra space below each item

# total height one item occupies (number + sep + gap + title + subtitle + pad)
ITEM_H = NUMBER_H + SEP_H + GAP_SEP_TITLE + TITLE_H + SUBTITLE_H + ITEM_BOTTOM_PAD


# ── helpers ───────────────────────────────────────────────────────────────

def _hex_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _textbox(
    slide,
    text: str,
    x: float, y: float, w: float, h: float,
    size: int,
    color: str,
    font: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    """Add a single-run textbox at the given position (all values in inches)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    if not p.runs:
        return
    r = p.runs[0]
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.color.rgb = _hex_rgb(color)
    r.font.name      = font


def _rect(slide, x: float, y: float, w: float, h: float, color: str) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_rgb(color)
    shape.line.fill.background()


# ── grid layout calculator ────────────────────────────────────────────────

def _grid(n: int) -> tuple[int, int]:
    """Return (cols, rows) for n items."""
    if n <= 2:
        return n, 1
    if n <= 4:
        return 2, math.ceil(n / 2)
    if n <= 6:
        return 2, math.ceil(n / 2)
    return 3, math.ceil(n / 3)


# ── main render function ──────────────────────────────────────────────────

def render_agenda(pptx_slide, slide: SlideData, theme: Theme) -> None:
    """
    Draw the AGENDA composition onto *pptx_slide*.

    Expects slide.content["items"] to be a list of dicts with keys:
        number   : str  — e.g. "01", "02" (AI-generated or auto-assigned)
        title    : str  — short item heading
        subtitle : str  — one-line description (optional, falls back to "")
    """
    t = theme
    items: list[dict[str, Any]] = slide.content.get("items", [])

    if not items:
        logger.warning("AGENDA slide index=%d has no items — rendering empty", slide.index)

    # Auto-assign numbers if missing
    for i, item in enumerate(items):
        if not item.get("number"):
            item["number"] = f"{i + 1:02d}"

    # ── slide title + accent rule ─────────────────────────────────────────
    title_size = 30 if len(slide.title) <= 40 else 24
    _textbox(
        pptx_slide, slide.title,
        TK.margin_outer, 0.20,
        TK.content_width, TK.title_height,
        title_size, t.primary, t.font_heading, bold=True,
    )
    _rect(
        pptx_slide,
        TK.margin_outer,
        0.20 + TK.title_height + TK.gap_title_rule,
        TK.content_width,
        TK.title_rule_height,
        t.accent,
    )

    # ── content area bounds ───────────────────────────────────────────────
    content_y = TK.content_start_y + 0.05   # a bit of breathing room
    content_h = SLIDE_H - content_y - 0.25  # bottom margin
    content_w = TK.content_width

    n = len(items)
    if n == 0:
        return

    cols, rows = _grid(n)

    # Column width with inter-column gap
    col_gap  = 0.35
    col_w    = (content_w - col_gap * (cols - 1)) / cols

    # Row height: distribute available vertical space, capped at ITEM_H
    row_gap  = 0.25
    row_h    = min(ITEM_H, (content_h - row_gap * (rows - 1)) / rows)

    # Recompute sub-element heights proportionally if row_h < ITEM_H
    scale    = row_h / ITEM_H
    num_h    = NUMBER_H * scale
    sep_h    = max(SEP_H, 0.025)
    gap_st   = GAP_SEP_TITLE * scale
    ttl_h    = TITLE_H * scale
    sub_h    = SUBTITLE_H * scale

    # Actual number font size (scale down for tight grids)
    num_size = max(32, int(NUMBER_SIZE * scale))
    ttl_size = max(11, int(TITLE_SIZE * scale))
    sub_size = max(9,  int(SUBTITLE_SIZE * scale))

    for idx, item in enumerate(items):
        col = idx % cols
        row = idx // cols

        ix = TK.margin_outer + col * (col_w + col_gap)
        iy = content_y       + row * (row_h + row_gap)

        number   = str(item.get("number", f"{idx + 1:02d}"))
        title    = str(item.get("title",    ""))
        subtitle = str(item.get("subtitle", ""))

        # ── ordinal number ────────────────────────────────────────────────
        _textbox(
            pptx_slide, number,
            ix, iy, col_w, num_h,
            num_size, t.accent, t.font_heading, bold=True,
        )

        # ── thin separator line under number ─────────────────────────────
        sep_y = iy + num_h
        _rect(pptx_slide, ix, sep_y, col_w * 0.65, sep_h, t.accent)

        # ── item title ────────────────────────────────────────────────────
        ttl_y = sep_y + sep_h + gap_st
        _textbox(
            pptx_slide, title,
            ix, ttl_y, col_w, ttl_h,
            ttl_size, t.primary, t.font_heading, bold=True,
        )

        # ── item subtitle ─────────────────────────────────────────────────
        if subtitle:
            sub_y = ttl_y + ttl_h
            _textbox(
                pptx_slide, subtitle,
                ix, sub_y, col_w, sub_h,
                sub_size, t.text_dark, t.font_body,
            )

    logger.debug(
        "AGENDA rendered: %d items, %dx%d grid, row_h=%.2f scale=%.2f",
        n, cols, rows, row_h, scale,
    )
