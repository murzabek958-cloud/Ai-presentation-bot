"""
Tests for PPTXBuilder.
No external API calls — all tests run fully offline.
"""
from datetime import datetime
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from ai.schemas import PlanMetadata, SlideData, SlideLayout
from presentation.builder import BuilderError, PPTXBuilder
from presentation.styles import get_theme

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

THEME = get_theme("academic")


def _slide(layout: SlideLayout, content: dict, **kwargs) -> SlideData:
    """Build a minimal valid SlideData for a given layout."""
    return SlideData(
        index=0,
        layout=layout,
        title="Test Slide",
        content=content,
        **kwargs,
    )


# Minimal content for each layout
LAYOUT_CONTENTS: dict[SlideLayout, dict] = {
    SlideLayout.TITLE: {"subtitle": "A subtitle", "author": "Author Name"},
    SlideLayout.TITLE_TEXT: {"body": "Some body text here."},
    SlideLayout.IMAGE_TEXT: {"body": "Text beside image."},
    SlideLayout.TWO_COLUMNS: {
        "left_title": "Left",
        "left_body": "Left column body.",
        "right_title": "Right",
        "right_body": "Right column body.",
    },
    SlideLayout.COMPARISON: {
        "left_label": "Option A",
        "left_points": ["Fast", "Cheap"],
        "right_label": "Option B",
        "right_points": ["Reliable", "Scalable"],
    },
    SlideLayout.TIMELINE: {
        "events": [
            {"year": "2020", "description": "First event"},
            {"year": "2022", "description": "Second event"},
        ]
    },
    SlideLayout.STATISTICS: {
        "stats": [
            {"value": "95%", "label": "Accuracy"},
            {"value": "1.2M", "label": "Users"},
        ]
    },
    SlideLayout.CHART: {
        "chart_type": "bar",
        "description": "Revenue by quarter.",
        "data_hint": "Q1: 10, Q2: 20, Q3: 15",
    },
    SlideLayout.QUOTE: {
        "quote": "The best way to predict the future is to create it.",
        "author": "Peter Drucker",
    },
    SlideLayout.CONCLUSION: {
        "summary": "We covered the key points.",
        "call_to_action": "Start building today.",
    },
}

# ---------------------------------------------------------------------------
# Presentation initialisation
# ---------------------------------------------------------------------------

def test_creates_16x9_presentation():
    builder = PPTXBuilder(theme=THEME)
    # Access internal prs to check dimensions
    prs: Presentation = builder._prs
    assert prs is not None
    assert prs.slide_width == Inches(13.33)
    assert prs.slide_height == Inches(7.5)


def test_create_presentation_resets_slides():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.TITLE, LAYOUT_CONTENTS[SlideLayout.TITLE]))
    assert len(builder._prs.slides) == 1

    builder.create_presentation()
    assert len(builder._prs.slides) == 0


# ---------------------------------------------------------------------------
# All 10 layouts — smoke tests (no crash)
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("layout", list(SlideLayout))
def test_add_slide_no_crash(layout: SlideLayout):
    """Every layout must render without raising an exception."""
    builder = PPTXBuilder(theme=THEME)
    slide = _slide(layout, LAYOUT_CONTENTS[layout])
    builder.add_slide(slide)  # must not raise
    assert len(builder._prs.slides) == 1


@pytest.mark.parametrize("layout", list(SlideLayout))
def test_add_slide_increments_count(layout: SlideLayout):
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(layout, LAYOUT_CONTENTS[layout]))
    builder.add_slide(_slide(layout, LAYOUT_CONTENTS[layout]))
    assert len(builder._prs.slides) == 2


# ---------------------------------------------------------------------------
# All 3 themes — smoke test
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("theme_name", ["academic", "modern", "minimal"])
def test_all_themes_render_title(theme_name: str):
    theme = get_theme(theme_name)
    builder = PPTXBuilder(theme=theme)
    builder.add_slide(_slide(SlideLayout.TITLE, LAYOUT_CONTENTS[SlideLayout.TITLE]))


# ---------------------------------------------------------------------------
# Missing / empty optional fields — no crash
# ---------------------------------------------------------------------------

def test_title_missing_author():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.TITLE, {"subtitle": "Sub only"}))


def test_title_empty_content():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.TITLE, {}))


def test_title_text_missing_body():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.TITLE_TEXT, {}))


def test_image_text_no_image_query():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.IMAGE_TEXT, {"body": "Text only."}))


def test_image_text_query_in_slide_field():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(
        _slide(SlideLayout.IMAGE_TEXT, {"body": "Body."}, image_query="sunset over mountains")
    )


def test_conclusion_no_cta():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.CONCLUSION, {"summary": "Done."}))


def test_quote_no_source():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.QUOTE, {"quote": "Words.", "author": "Someone"}))


def test_timeline_empty_events():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.TIMELINE, {"events": []}))


def test_statistics_empty_stats():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.STATISTICS, {"stats": []}))


def test_comparison_empty_points():
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.COMPARISON, {
        "left_label": "A", "left_points": [],
        "right_label": "B", "right_points": [],
    }))


# ---------------------------------------------------------------------------
# save() — file is created correctly
# ---------------------------------------------------------------------------

def test_save_creates_pptx_file(tmp_path: Path):
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.TITLE, LAYOUT_CONTENTS[SlideLayout.TITLE]))

    out = tmp_path / "output.pptx"
    result = builder.save(str(out))

    assert Path(result).exists()
    assert Path(result).suffix == ".pptx"
    assert Path(result).stat().st_size > 0


def test_save_creates_parent_dirs(tmp_path: Path):
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.TITLE_TEXT, LAYOUT_CONTENTS[SlideLayout.TITLE_TEXT]))

    nested = tmp_path / "a" / "b" / "c" / "out.pptx"
    builder.save(str(nested))
    assert nested.exists()


def test_save_returns_path_string(tmp_path: Path):
    builder = PPTXBuilder(theme=THEME)
    builder.add_slide(_slide(SlideLayout.CONCLUSION, LAYOUT_CONTENTS[SlideLayout.CONCLUSION]))

    result = builder.save(str(tmp_path / "deck.pptx"))
    assert isinstance(result, str)


def test_save_all_layouts(tmp_path: Path):
    """One slide per layout — full deck saved successfully."""
    builder = PPTXBuilder(theme=THEME)
    for i, (layout, content) in enumerate(LAYOUT_CONTENTS.items()):
        slide = SlideData(index=i, layout=layout, title=f"Slide {i}", content=content)
        builder.add_slide(slide)

    out = tmp_path / "full_deck.pptx"
    builder.save(str(out))
    assert out.exists()
    assert out.stat().st_size > 0


# ---------------------------------------------------------------------------
# Error handling
# ---------------------------------------------------------------------------

def test_save_raises_if_not_initialised():
    builder = PPTXBuilder(theme=THEME)
    builder._prs = None  # simulate uninitialised state
    with pytest.raises(BuilderError):
        builder.save("/tmp/should_not_exist.pptx")


def test_add_slide_raises_if_not_initialised():
    builder = PPTXBuilder(theme=THEME)
    builder._prs = None
    with pytest.raises(BuilderError):
        builder.add_slide(_slide(SlideLayout.TITLE, {}))
