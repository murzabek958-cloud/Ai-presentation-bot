"""
Integration tests for topic-aware theme selection in PresentationRenderer.

Verifies the complete flow:
    plan.topic → ThemeSelector → Theme → PresentationRenderer → PPTXBuilder

All tests are offline — no external API calls.
"""
from __future__ import annotations

import re
from datetime import datetime

import pytest
from pptx import Presentation

from ai.schemas import PlanMetadata, PresentationPlan, SlideData, SlideLayout
from presentation.renderer import PresentationRenderer, RendererError
from presentation.styles import get_theme
from presentation.theme_selector import ThemeSelector

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_selector = ThemeSelector()


def _make_plan(
    topic: str,
    style: str = "academic",
    layouts: list[SlideLayout] | None = None,
) -> PresentationPlan:
    if layouts is None:
        layouts = [SlideLayout.TITLE, SlideLayout.TITLE_TEXT, SlideLayout.CONCLUSION]

    content_for = {
        SlideLayout.TITLE:       {"subtitle": "Test subtitle", "author": "Test"},
        SlideLayout.TITLE_TEXT:  {"body": "Test body text for the slide."},
        SlideLayout.CONCLUSION:  {"summary": "Test summary.", "call_to_action": ""},
        SlideLayout.STATISTICS:  {"stats": [{"value": "99%", "label": "Test"}]},
        SlideLayout.TWO_COLUMNS: {"left_title": "L", "left_body": "Left.", "right_title": "R", "right_body": "Right."},
        SlideLayout.QUOTE:       {"quote": "Test quote.", "author": "Author"},
        SlideLayout.TIMELINE:    {"events": [{"year": "2024", "description": "Event"}]},
        SlideLayout.COMPARISON:  {"left_label": "A", "left_points": ["x"], "right_label": "B", "right_points": ["y"]},
        SlideLayout.IMAGE_TEXT:  {"body": "Image body."},
        SlideLayout.CHART:       {"chart_type": "bar", "description": "desc"},
        SlideLayout.AGENDA:      {"items": [
            {"number": "01", "title": "Topic One", "subtitle": "First section"},
            {"number": "02", "title": "Topic Two", "subtitle": "Second section"},
        ]},
    }

    slides = [
        SlideData(
            index=i,
            layout=layout,
            title=f"Slide {i}",
            content=content_for.get(layout, {}),
        )
        for i, layout in enumerate(layouts)
    ]
    return PresentationPlan(
        topic=topic,
        style=style,
        slide_count=len(slides),
        slides=slides,
        metadata=PlanMetadata(language="kk", generated_at=datetime.utcnow()),
    )


def _hex(value: str) -> bool:
    return bool(re.fullmatch(r"#[0-9A-Fa-f]{6}", value))


def _render_bytes(plan: PresentationPlan, explicit: bool) -> bytes:
    """Render plan to in-memory bytes via a temp file."""
    import tempfile, os
    renderer = PresentationRenderer(plan, style_is_explicit=explicit)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    try:
        renderer.save(path)
        with open(path, "rb") as f:
            return f.read()
    finally:
        os.unlink(path)


# ---------------------------------------------------------------------------
# 1. Topic-aware palette selection (style_is_explicit=False)
# ---------------------------------------------------------------------------

class TestTopicAwarePalette:

    def test_anatomy_uses_medicine_palette(self):
        plan = _make_plan("Адам анатомиясы")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        profile = _selector.profile(plan.topic)
        # Renderer theme must match what ThemeSelector picks for the topic
        expected = _selector.select(plan.topic)
        assert renderer.theme.primary == expected.primary
        assert renderer.theme.accent == expected.accent
        assert profile.primary_category == "medicine"

    def test_caspian_uses_ocean_ecology_palette(self):
        plan = _make_plan("Каспий теңізінің экологиясы")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        profile = _selector.profile(plan.topic)
        expected = _selector.select(plan.topic)
        assert renderer.theme.primary == expected.primary
        assert renderer.theme.accent == expected.accent
        assert profile.palette_name in ("ocean_ecology", "ocean", "nature")

    def test_ai_uses_technology_palette(self):
        plan = _make_plan("Жасанды интеллект")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        profile = _selector.profile(plan.topic)
        assert profile.primary_category == "technology"
        expected = _selector.select(plan.topic)
        assert renderer.theme.primary == expected.primary

    def test_history_uses_history_palette(self):
        plan = _make_plan("Қазақстан тарихы")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        profile = _selector.profile(plan.topic)
        assert profile.primary_category == "history"
        expected = _selector.select(plan.topic)
        assert renderer.theme.primary == expected.primary

    def test_business_uses_business_palette(self):
        plan = _make_plan("Кәсіпкерлік")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        profile = _selector.profile(plan.topic)
        assert profile.primary_category == "business"

    def test_unknown_topic_uses_neutral_theme(self):
        plan = _make_plan("абракадабра xyz123 непонятное")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        profile = _selector.profile(plan.topic)
        assert profile.palette_name == "neutral"
        expected = _selector.select(plan.topic)
        assert renderer.theme.primary == expected.primary


# ---------------------------------------------------------------------------
# 2. Explicit style must NOT be overridden by topic palette
# ---------------------------------------------------------------------------

class TestExplicitStylePreserved:

    @pytest.mark.parametrize("explicit_style", ["academic", "modern", "minimal"])
    def test_explicit_style_preserved(self, explicit_style: str):
        """When style_is_explicit=True, renderer must use the named style Theme."""
        plan = _make_plan("Адам анатомиясы", style=explicit_style)
        renderer = PresentationRenderer(plan, style_is_explicit=True)
        expected_theme = get_theme(explicit_style)
        assert renderer.theme.primary == expected_theme.primary
        assert renderer.theme.accent == expected_theme.accent
        assert renderer.theme.name == expected_theme.name

    def test_explicit_academic_not_overridden_by_medicine(self):
        """anatomy topic + explicit academic → must NOT use medicine palette."""
        plan = _make_plan("Адам анатомиясы", style="academic")
        renderer_explicit = PresentationRenderer(plan, style_is_explicit=True)
        renderer_auto    = PresentationRenderer(plan, style_is_explicit=False)
        # Explicit must match academic theme exactly
        assert renderer_explicit.theme.primary == get_theme("academic").primary
        # Auto must differ from explicit (medicine palette ≠ academic)
        # (This also confirms that auto-selection actually does something)
        assert renderer_auto.theme.primary != renderer_explicit.theme.primary

    def test_explicit_modern_not_overridden_by_history(self):
        plan = _make_plan("Қазақстан тарихы", style="modern")
        renderer = PresentationRenderer(plan, style_is_explicit=True)
        assert renderer.theme.primary == get_theme("modern").primary

    def test_style_is_explicit_false_is_default(self):
        """Omitting style_is_explicit defaults to False (topic-aware mode)."""
        plan_anatomy = _make_plan("Адам анатомиясы")
        r_default  = PresentationRenderer(plan_anatomy)           # no kwarg
        r_explicit_false = PresentationRenderer(plan_anatomy, style_is_explicit=False)
        assert r_default.theme.primary == r_explicit_false.theme.primary


# ---------------------------------------------------------------------------
# 3. theme property is accessible before and after render
# ---------------------------------------------------------------------------

class TestThemeProperty:

    def test_theme_property_before_render(self):
        plan = _make_plan("Жасанды интеллект")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        assert _hex(renderer.theme.primary)
        assert _hex(renderer.theme.accent)

    def test_theme_property_is_theme_instance(self):
        from presentation.styles import Theme
        plan = _make_plan("Каспий теңізінің экологиясы")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        assert isinstance(renderer.theme, Theme)


# ---------------------------------------------------------------------------
# 4. Rendered PPTX is valid and uses correct colors
# ---------------------------------------------------------------------------

class TestRenderedPptxColors:
    """
    Verify that the selected palette colors actually appear in the generated
    .pptx file, not just in the theme object.
    """

    def _primary_rgb_from_hex(self, hex_color: str) -> tuple[int, int, int]:
        h = hex_color.lstrip("#")
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

    def _collect_slide_colors(self, pptx_bytes: bytes) -> set[str]:
        """Extract all solid fill hex colors from all shapes in all slides."""
        import io
        prs = Presentation(io.BytesIO(pptx_bytes))
        colors: set[str] = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    fill = shape.fill
                    if fill.type == 1:  # MSO_FILL.SOLID
                        rgb = fill.fore_color.rgb
                        colors.add(f"#{rgb}".upper())
                except Exception:
                    pass
        return colors

    def test_anatomy_pptx_contains_medicine_primary_color(self, tmp_path):
        plan = _make_plan("Адам анатомиясы")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        out = str(tmp_path / "anatomy.pptx")
        renderer.save(out)
        with open(out, "rb") as f:
            pptx_bytes = f.read()
        colors = self._collect_slide_colors(pptx_bytes)
        expected_primary = renderer.theme.primary.upper()
        assert expected_primary in colors, (
            f"Expected primary color {expected_primary} in PPTX shapes, got: {colors}"
        )

    def test_history_pptx_contains_history_primary_color(self, tmp_path):
        plan = _make_plan("Қазақстан тарихы")
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        out = str(tmp_path / "history.pptx")
        renderer.save(out)
        with open(out, "rb") as f:
            pptx_bytes = f.read()
        colors = self._collect_slide_colors(pptx_bytes)
        expected_primary = renderer.theme.primary.upper()
        assert expected_primary in colors, (
            f"Expected primary color {expected_primary} in PPTX shapes, got: {colors}"
        )

    def test_explicit_academic_pptx_uses_academic_colors(self, tmp_path):
        plan = _make_plan("Адам анатомиясы", style="academic")
        renderer = PresentationRenderer(plan, style_is_explicit=True)
        out = str(tmp_path / "academic.pptx")
        renderer.save(out)
        with open(out, "rb") as f:
            pptx_bytes = f.read()
        colors = self._collect_slide_colors(pptx_bytes)
        academic_primary = get_theme("academic").primary.upper()
        assert academic_primary in colors


# ---------------------------------------------------------------------------
# 5. Existing-style presentations still work exactly as before
# ---------------------------------------------------------------------------

class TestBackwardCompatibility:

    @pytest.mark.parametrize("style", ["academic", "modern", "minimal"])
    def test_existing_style_renders_without_error(self, style: str, tmp_path):
        plan = _make_plan("Test topic", style=style)
        renderer = PresentationRenderer(plan, style_is_explicit=True)
        out = str(tmp_path / f"{style}.pptx")
        renderer.save(out)
        assert (tmp_path / f"{style}.pptx").exists()

    def test_renderer_without_kwarg_still_works(self, tmp_path):
        """Existing call sites: PresentationRenderer(plan) — no crash."""
        plan = _make_plan("Some topic")
        renderer = PresentationRenderer(plan)
        out = str(tmp_path / "compat.pptx")
        renderer.save(out)
        assert (tmp_path / "compat.pptx").exists()

    def test_all_layouts_render_with_topic_theme(self, tmp_path):
        layouts = list(SlideLayout)
        plan = _make_plan("Жасанды интеллект", layouts=layouts)
        renderer = PresentationRenderer(plan, style_is_explicit=False)
        out = str(tmp_path / "all_layouts.pptx")
        renderer.save(out)
        assert (tmp_path / "all_layouts.pptx").exists()
